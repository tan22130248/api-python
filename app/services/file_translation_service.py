import logging
import io
from typing import Dict, Any
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

logger = logging.getLogger(__name__)


def translate_docx_file(content: bytes, source_lang: str = "vi", target_lang: str = "en") -> bytes:
    """
    Translate .docx file while preserving structure.
    Only translates paragraph text, keeps formatting.
    """
    from app.services.translation_service import translate_text

    doc = Document(io.BytesIO(content))

    for para in doc.paragraphs:
        if para.text.strip():
            try:
                translated = translate_text(para.text, source_lang, target_lang)
                # Clear existing runs and add translated text
                for run in para.runs:
                    run.text = ""
                if para.runs:
                    para.runs[0].text = translated
                else:
                    para.add_run(translated)
            except Exception as e:
                logger.warning(f"Failed to translate paragraph: {e}")

    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()


def translate_txt_file(content: bytes, source_lang: str = "vi", target_lang: str = "en") -> bytes:
    """
    Translate plain text file.
    """
    from app.services.translation_service import translate_document

    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        text = content.decode("cp1252", errors="ignore")

    result = translate_document(text, source_lang, target_lang)
    translated_text = result["translated_text"]

    return translated_text.encode("utf-8")


def translate_pptx_file(content: bytes, source_lang: str = "vi", target_lang: str = "en") -> bytes:
    """
    Translate .pptx file while preserving structure.
    Only translates text in shapes, keeps formatting.
    """
    from pptx import Presentation
    from app.services.translation_service import translate_text

    prs = Presentation(io.BytesIO(content))

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                try:
                    translated = translate_text(shape.text, source_lang, target_lang)
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ""
                        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                            shape.text_frame.paragraphs[0].runs[0].text = translated
                        else:
                            shape.text_frame.paragraphs[0].text = translated
                except Exception as e:
                    logger.warning(f"Failed to translate shape text: {e}")

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
