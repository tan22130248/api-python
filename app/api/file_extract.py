from fastapi import APIRouter, HTTPException, UploadFile, File
import logging
import io

router = APIRouter(prefix="/api/translate", tags=["translate"])
logger = logging.getLogger(__name__)

@router.post("/extract")
async def extract_text_from_file(file: UploadFile = File(...)):
    """
    Extract text from an uploaded document.
    Supported formats: .txt, .pdf, .docx, .xlsx, .pptx
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="Tên tệp không hợp lệ")

    filename = file.filename.lower()
    content = await file.read()
    extracted_text = ""

    try:
        if filename.endswith(".txt"):
            try:
                extracted_text = content.decode("utf-8")
            except UnicodeDecodeError:
                # Fallback for Windows-1258 / CP1252 if UTF-8 fails
                extracted_text = content.decode("cp1252", errors="ignore")

        elif filename.endswith(".pdf"):
            import pypdf
            pdf_reader = pypdf.PdfReader(io.BytesIO(content))
            text_parts = []
            for page in pdf_reader.pages:
                text_parts.append(page.extract_text() or "")
            extracted_text = "\n\n".join(text_parts)

        elif filename.endswith(".docx"):
            from docx import Document
            doc = Document(io.BytesIO(content))
            text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
            extracted_text = "\n".join(text_parts)

        elif filename.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            extracted_text = "\n\n".join(text_parts)

        elif filename.endswith(".xlsx"):
            from openpyxl import load_workbook
            wb = load_workbook(filename=io.BytesIO(content), data_only=True)
            text_parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text_parts.append(row_text)
            extracted_text = "\n".join(text_parts)

        else:
            raise HTTPException(
                status_code=400, 
                detail="Định dạng tệp không được hỗ trợ. Vui lòng tải lên .txt, .pdf, .docx, .xlsx, hoặc .pptx"
            )

        extracted_text = extracted_text.strip()
        if not extracted_text:
            raise HTTPException(status_code=400, detail="Không tìm thấy văn bản trong tệp")

        # Trim text to max 50000 characters if it's too long
        if len(extracted_text) > 50000:
            extracted_text = extracted_text[:50000]

        return {
            "success": True,
            "message": "Trích xuất văn bản thành công",
            "data": {
                "text": extracted_text,
                "filename": file.filename
            }
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error extracting text from {file.filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Lỗi trích xuất văn bản: {str(e)}")
